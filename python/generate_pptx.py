#!/usr/bin/env python3
"""PPTX proposal generator. Reads JSON from stdin, writes .pptx file."""
import sys
import json
from datetime import datetime, timedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION


def hex_to_rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class ProposalGenerator:
    def __init__(self, analysis, config, output_path):
        self.a = analysis
        self.c = config
        self.output_path = output_path
        self.prs = Presentation()
        self.prs.slide_width = Inches(config['slide_size']['width'])
        self.prs.slide_height = Inches(config['slide_size']['height'])

        colors = config['colors']
        self.PRIMARY = hex_to_rgb(colors['primary'])
        self.DARK = hex_to_rgb(colors['dark'])
        self.GRAY = hex_to_rgb(colors['gray'])
        self.LGRAY = hex_to_rgb(colors['light_gray'])
        self.WHITE = hex_to_rgb(colors['white'])
        self.BG = hex_to_rgb(colors['bg_light'])
        self.BG_WARM = hex_to_rgb(colors['bg_warm'])
        self.FONT = config['fonts']['primary']

    def tb(self, slide, l, t, w, h, text='', sz=14, bold=False, color=None, align=PP_ALIGN.LEFT):
        color = color or self.DARK
        box = slide.shapes.add_textbox(l, t, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = self.FONT
        p.alignment = align
        return tf

    def bullets(self, slide, l, t, w, h, items, sz=14, color=None, sp=Pt(6)):
        color = color or self.DARK
        box = slide.shapes.add_textbox(l, t, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(item)
            p.font.size = Pt(sz)
            p.font.color.rgb = color
            p.font.name = self.FONT
            p.space_after = sp
        return tf

    def line(self, slide, l, t, w):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(3))
        s.fill.solid()
        s.fill.fore_color.rgb = self.PRIMARY
        s.line.fill.background()

    def box(self, slide, l, t, w, h, color=None):
        color = color or self.BG
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        s.line.width = Pt(1)
        s.shadow.inherit = False
        return s

    def flow_box(self, slide, l, t, w, h, text, fill_color=None, sz=11):
        fill_color = fill_color or self.PRIMARY
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = True
        p.font.color.rgb = self.WHITE
        p.font.name = self.FONT
        p.alignment = PP_ALIGN.CENTER
        return shape

    def arrow_right(self, slide, l, t, w, h):
        s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        s.line.fill.background()
        s.shadow.inherit = False

    def arrow_down(self, slide, l, t, w, h):
        s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        s.line.fill.background()
        s.shadow.inherit = False

    def new_slide(self):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = self.WHITE
        return s

    def slide_title(self, slide, text, num=None):
        label = f'{num}. {text}' if num else text
        self.tb(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6), label, sz=28, bold=True)
        self.line(slide, Inches(0.8), Inches(1.0), Inches(11.5))

    def build_cover(self):
        s = self.new_slide()
        self.tb(s, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
                self.a['title'], sz=38, bold=True, align=PP_ALIGN.CENTER)
        self.tb(s, Inches(1.5), Inches(2.9), Inches(10), Inches(0.6),
                self.a['subtitle'], sz=18, color=self.GRAY, align=PP_ALIGN.CENTER)
        self.line(s, Inches(5), Inches(3.6), Inches(3.333))

        today = datetime.now()
        validity = today + timedelta(days=self.c['proposal'].get('validity_days', 30))
        company = self.c.get('company', {})
        proposer = company.get('name', '')
        if company.get('title'):
            proposer = f"{proposer} / {company['title']}" if proposer else company['title']

        meta = [
            ('提案日期', today.strftime('%Y 年 %m 月 %d 日')),
            ('有效期限', validity.strftime('%Y 年 %m 月 %d 日')),
        ]
        if proposer:
            meta.append(('提案人', proposer))

        for i, (label, val) in enumerate(meta):
            y = Inches(4.2) + Inches(i * 0.45)
            self.tb(s, Inches(4), y, Inches(2), Inches(0.4), label, sz=13, color=self.LGRAY, align=PP_ALIGN.RIGHT)
            self.tb(s, Inches(6.3), y, Inches(4), Inches(0.4), val, sz=13)

    def build_background(self):
        s = self.new_slide()
        self.slide_title(s, '專案背景與目標', 1)
        self.tb(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.5),
                self.a['background'], sz=15, color=self.GRAY)
        self.tb(s, Inches(0.8), Inches(3.2), Inches(4), Inches(0.5), '核心目標', sz=20, bold=True)
        self.bullets(s, Inches(1.0), Inches(3.8), Inches(11), Inches(3.5),
                     self.a.get('goals', []), sz=15, sp=Pt(10))

    def build_architecture(self):
        arch = self.a.get('architecture')
        if not arch:
            return
        s = self.new_slide()
        self.slide_title(s, '系統架構', 2)

        # Visual flow with colored boxes and arrows
        tech_stack = arch.get('tech_stack', [])
        COLORS = [
            self.PRIMARY,
            RGBColor(0x34, 0x98, 0xDB),
            RGBColor(0x2E, 0xCC, 0x71),
            RGBColor(0x9B, 0x59, 0xB6),
            RGBColor(0xE7, 0x4C, 0x3C),
            RGBColor(0x1A, 0xBC, 0x9C),
        ]

        BOX_W = Inches(2.2)
        BOX_H = Inches(0.7)
        ARR_W = Inches(0.5)
        ARR_H = Inches(0.3)

        # Render up to 4 boxes in a row with arrows
        top_items = tech_stack[:4]
        row1_y = Inches(1.5)
        x_positions = []
        for i, tech in enumerate(top_items):
            x = Inches(0.8 + i * 3.0)
            x_positions.append(x)
            self.flow_box(s, x, row1_y, BOX_W, BOX_H, tech['name'], fill_color=COLORS[i % len(COLORS)])
            if i < len(top_items) - 1:
                self.arrow_right(s, x + BOX_W + Inches(0.1), row1_y + Inches(0.2), ARR_W, ARR_H)

        # If more than 4 items, add second row
        bottom_items = tech_stack[4:]
        if bottom_items:
            if x_positions:
                mid_x = x_positions[len(x_positions) // 2]
                self.arrow_down(s, mid_x + Inches(0.95), Inches(2.3), Inches(0.3), Inches(0.5))
            row2_y = Inches(3.0)
            for i, tech in enumerate(bottom_items):
                x = Inches(0.8 + i * 3.0)
                self.flow_box(s, x, row2_y, BOX_W, BOX_H, tech['name'], fill_color=COLORS[(i + 4) % len(COLORS)])
                if i < len(bottom_items) - 1:
                    self.arrow_right(s, x + BOX_W + Inches(0.1), row2_y + Inches(0.2), ARR_W, ARR_H)

        # Tech descriptions
        desc_y = Inches(4.2) if bottom_items else Inches(2.8)
        self.tb(s, Inches(0.8), desc_y, Inches(4), Inches(0.5), '技術選型', sz=20, bold=True)
        tech_items = [f"{t['name']}: {t['description']}" for t in tech_stack]
        if tech_items:
            self.bullets(s, Inches(1.0), desc_y + Inches(0.6), Inches(11), Inches(3),
                         tech_items, sz=14, sp=Pt(8))

    def build_pricing_comparison(self):
        """Bar chart comparing plans side by side."""
        plans = self.a.get('pricing', {}).get('plans', [])
        if len(plans) < 2:
            return
        s = self.new_slide()
        self.slide_title(s, '方案對比')

        plan_a = plans[0]
        plan_b = plans[1]

        # Build chart data
        categories = [item['name'] for item in plan_b.get('items', [])]
        values_a = []
        values_b = [item.get('price', 0) for item in plan_b.get('items', [])]

        a_items = {item['name']: item.get('price', 0) for item in plan_a.get('items', [])}
        for cat in categories:
            values_a.append(a_items.get(cat, 0))

        chart_data = CategoryChartData()
        chart_data.categories = categories
        currency = plan_a.get('currency', 'NT$')
        chart_data.add_series(f"{plan_a['label']}: {currency} {plan_a['total']:,}", values_a)
        chart_data.add_series(f"{plan_b['label']}: {currency} {plan_b['total']:,}", values_b)

        chart_frame = s.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.8), Inches(1.3), Inches(7.5), Inches(5.5),
            chart_data
        )
        chart = chart_frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(12)
        chart.legend.font.name = self.FONT

        plot = chart.plots[0]
        plot.gap_width = 120
        plot.series[0].format.fill.solid()
        plot.series[0].format.fill.fore_color.rgb = RGBColor(0xBD, 0xBD, 0xBD)
        plot.series[1].format.fill.solid()
        plot.series[1].format.fill.fore_color.rgb = self.PRIMARY

        value_axis = chart.value_axis
        value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        value_axis.tick_labels.font.size = Pt(10)
        value_axis.tick_labels.font.name = self.FONT
        value_axis.tick_labels.number_format = '#,##0'

        cat_axis = chart.category_axis
        cat_axis.tick_labels.font.size = Pt(11)
        cat_axis.tick_labels.font.name = self.FONT

        # Right side: differences
        b_only = [item['name'] for item in plan_b.get('items', []) if item['name'] not in a_items or a_items[item['name']] == 0]
        b_upgraded = [item['name'] for item in plan_b.get('items', []) if item['name'] in a_items and a_items[item['name']] > 0 and item.get('price', 0) > a_items[item['name']]]

        highlights = b_only + b_upgraded
        if highlights:
            self.tb(s, Inches(8.8), Inches(1.3), Inches(4), Inches(0.5),
                    f"{plan_b['label']} 額外包含", sz=16, bold=True)
            self.bullets(s, Inches(9.0), Inches(1.9), Inches(3.5), Inches(3),
                         highlights, sz=13, sp=Pt(8))

        diff = plan_b['total'] - plan_a['total']
        if diff > 0 and b_only:
            self.box(s, Inches(8.8), Inches(4.8), Inches(3.8), Inches(1.5), color=self.BG_WARM)
            self.tb(s, Inches(9.0), Inches(4.95), Inches(3.5), Inches(0.4),
                    f"差額僅 {currency} {diff:,}", sz=16, bold=True, color=self.PRIMARY, align=PP_ALIGN.CENTER)
            per_item = diff // max(len(b_only), 1)
            self.tb(s, Inches(9.0), Inches(5.45), Inches(3.5), Inches(0.6),
                    f"多 {len(b_only)} 項功能\n平均每項僅 {currency} {per_item:,}",
                    sz=12, color=self.GRAY, align=PP_ALIGN.CENTER)

    def build_gantt_timeline(self):
        """Gantt-style visual timeline."""
        timeline = self.a.get('timeline', [])
        if not timeline:
            return
        s = self.new_slide()
        self.slide_title(s, '時程規劃')

        GANTT_COLORS = [
            RGBColor(0x95, 0xA5, 0xA6),
            self.PRIMARY,
            RGBColor(0x34, 0x98, 0xDB),
            RGBColor(0x9B, 0x59, 0xB6),
            RGBColor(0x2E, 0xCC, 0x71),
            RGBColor(0x1A, 0xBC, 0x9C),
        ]

        total_phases = len(timeline)
        total_weeks = max(total_phases + 2, 7)
        gantt_left = Inches(3.8)
        gantt_width = Inches(8.2)
        week_w = gantt_width / total_weeks

        # Week headers
        for w in range(total_weeks):
            x = gantt_left + week_w * w
            self.tb(s, x, Inches(1.3), week_w, Inches(0.3),
                    f"W{w+1}", sz=10, bold=True, color=self.LGRAY, align=PP_ALIGN.CENTER)
            grid = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.6), Pt(0.5), Inches(total_phases * 0.55 + 0.2))
            grid.fill.solid()
            grid.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
            grid.line.fill.background()

        # Bars: auto-assign start positions sequentially
        cumulative = 0
        for i, phase in enumerate(timeline):
            y = Inches(1.7) + Inches(i * 0.55)
            dur = phase.get('duration', '1 week')

            # Parse duration to weeks estimate
            dur_weeks = 1
            if '2' in dur and '3' in dur:
                dur_weeks = 3
            elif '1' in dur and '2' in dur:
                dur_weeks = 2
            elif '3' in dur or '4' in dur:
                dur_weeks = 1
            else:
                dur_weeks = 1

            # Label
            self.tb(s, Inches(0.8), y, Inches(2.8), Inches(0.45),
                    phase['phase'], sz=12, bold=True)

            # Bar
            x = gantt_left + week_w * cumulative
            w = week_w * dur_weeks
            bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y + Inches(0.05), w, Inches(0.35))
            bar.fill.solid()
            bar.fill.fore_color.rgb = GANTT_COLORS[i % len(GANTT_COLORS)]
            bar.line.fill.background()
            bar.shadow.inherit = False

            tf = bar.text_frame
            p = tf.paragraphs[0]
            p.text = dur
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = self.WHITE
            p.font.name = self.FONT
            p.alignment = PP_ALIGN.CENTER

            cumulative += dur_weeks

        ts = self.a.get('timeline_summary', '')
        if ts:
            self.tb(s, Inches(0.8), Inches(1.7) + Inches(total_phases * 0.55) + Inches(0.3),
                    Inches(11), Inches(0.3), ts, sz=12, color=self.LGRAY)

    def build_features(self):
        features = self.a.get('features', [])
        if not features:
            return
        s = self.new_slide()
        self.slide_title(s, '功能範圍', 3)

        # Layout: up to 2 columns top, up to 3 columns bottom
        cols = min(len(features), 2)
        col_w = 5.5
        top_features = features[:2]
        bottom_features = features[2:5]

        for i, feat in enumerate(top_features):
            x = Inches(0.8 + i * 6.0)
            self.tb(s, x, Inches(1.3), Inches(col_w), Inches(0.4),
                    feat['category'], sz=17, bold=True)
            self.bullets(s, x + Inches(0.2), Inches(1.8), Inches(col_w - 0.2), Inches(2.5),
                         feat.get('items', []), sz=13, sp=Pt(4))

        if bottom_features:
            bcol_w = 3.6
            for i, feat in enumerate(bottom_features):
                x = Inches(0.8 + i * 4.0)
                self.tb(s, x, Inches(4.5), Inches(bcol_w), Inches(0.4),
                        feat['category'], sz=17, bold=True)
                self.bullets(s, x + Inches(0.2), Inches(5.0), Inches(bcol_w - 0.2), Inches(2.2),
                             feat.get('items', []), sz=13, sp=Pt(4))

    def build_pricing(self, plan, slide_num, is_recommended=False):
        s = self.new_slide()
        if slide_num == 4:
            self.slide_title(s, '方案與報價', 4)

        bg_color = self.BG_WARM if is_recommended else self.BG
        b = self.box(s, Inches(0.8), Inches(0.4 if slide_num > 4 else 1.3),
                     Inches(11.5), Inches(6.2 if slide_num > 4 else 5.8), color=bg_color)
        if is_recommended:
            b.line.color.rgb = self.PRIMARY

        y_base = Inches(0.6 if slide_num > 4 else 1.5)
        self.tb(s, Inches(1.2), y_base, Inches(6), Inches(0.5),
                plan['label'], sz=22, bold=True)

        if is_recommended:
            badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(7.5), y_base + Inches(0.05), Inches(0.9), Inches(0.35))
            badge.fill.solid()
            badge.fill.fore_color.rgb = self.PRIMARY
            badge.line.fill.background()
            self.tb(s, Inches(7.5), y_base + Inches(0.05), Inches(0.9), Inches(0.35),
                    '推薦', sz=12, bold=True, color=self.WHITE, align=PP_ALIGN.CENTER)

        currency = plan.get('currency', 'NT$')
        total = f"{currency} {plan['total']:,}"
        self.tb(s, Inches(1.2), y_base + Inches(0.6), Inches(6), Inches(0.6),
                total, sz=36, bold=True, color=self.PRIMARY)

        # Table headers
        y0 = y_base + Inches(1.4)
        self.tb(s, Inches(1.2), y0, Inches(2.8), Inches(0.35), '項目', sz=12, bold=True, color=self.GRAY)
        self.tb(s, Inches(4.0), y0, Inches(5.5), Inches(0.35), '說明', sz=12, bold=True, color=self.GRAY)
        self.tb(s, Inches(9.8), y0, Inches(2), Inches(0.35), '費用', sz=12, bold=True, color=self.GRAY, align=PP_ALIGN.RIGHT)

        items = plan.get('items', [])
        row_h = min(0.4, 2.5 / max(len(items), 1))
        for i, item in enumerate(items):
            y = y0 + Inches(0.4) + Inches(i * row_h)
            self.tb(s, Inches(1.2), y, Inches(2.8), Inches(0.35), item['name'], sz=13)
            self.tb(s, Inches(4.0), y, Inches(5.5), Inches(0.35), item.get('description', ''), sz=12, color=self.GRAY)
            price_str = f"{currency} {item['price']:,}" if isinstance(item['price'], (int, float)) else str(item['price'])
            self.tb(s, Inches(9.8), y, Inches(2), Inches(0.35), price_str, sz=13, align=PP_ALIGN.RIGHT)

        # Total line
        yt = y0 + Inches(0.4) + Inches(len(items) * row_h) + Inches(0.1)
        self.line(s, Inches(1.2), yt, Inches(10.6))
        self.tb(s, Inches(1.2), yt + Inches(0.1), Inches(4), Inches(0.4), '合計', sz=16, bold=True)
        self.tb(s, Inches(9.8), yt + Inches(0.1), Inches(2), Inches(0.4),
                total, sz=16, bold=True, color=self.PRIMARY, align=PP_ALIGN.RIGHT)

        # Summary
        summary = plan.get('summary', '')
        fit = plan.get('fit', '')
        note = f"包含: {summary}" if summary else ''
        if fit:
            note += f"\n適合: {fit}" if note else f"適合: {fit}"
        if note:
            self.tb(s, Inches(1.2), yt + Inches(0.7), Inches(10.5), Inches(0.8),
                    note, sz=11, color=self.LGRAY)

    def build_maintenance_timeline(self):
        s = self.new_slide()
        # Maintenance
        self.tb(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5), '月維護方案', sz=22, bold=True)
        self.line(s, Inches(0.8), Inches(0.95), Inches(5))

        for i, m in enumerate(self.a.get('maintenance', [])):
            y = Inches(1.2) + Inches(i * 0.45)
            self.tb(s, Inches(1.0), y, Inches(2), Inches(0.35), m['name'], sz=14, bold=True)
            self.tb(s, Inches(3.2), y, Inches(6), Inches(0.35), m['description'], sz=13, color=self.GRAY)
            self.tb(s, Inches(9.8), y, Inches(2.2), Inches(0.35), m['price'], sz=14, color=self.PRIMARY, align=PP_ALIGN.RIGHT)

        maint_count = len(self.a.get('maintenance', []))
        maint_bottom = Inches(1.2) + Inches(maint_count * 0.45) + Inches(0.3)

        # Timeline
        self.tb(s, Inches(0.8), maint_bottom, Inches(5), Inches(0.5), '時程規劃', sz=22, bold=True)
        self.line(s, Inches(0.8), maint_bottom + Inches(0.55), Inches(11.5))

        y0 = maint_bottom + Inches(0.8)
        self.tb(s, Inches(1.0), y0, Inches(2.8), Inches(0.35), '階段', sz=12, bold=True, color=self.GRAY)
        self.tb(s, Inches(4.0), y0, Inches(5.5), Inches(0.35), '內容', sz=12, bold=True, color=self.GRAY)
        self.tb(s, Inches(10), y0, Inches(2), Inches(0.35), '預估時間', sz=12, bold=True, color=self.GRAY, align=PP_ALIGN.RIGHT)

        for i, t in enumerate(self.a.get('timeline', [])):
            y = y0 + Inches(0.45) + Inches(i * 0.42)
            self.tb(s, Inches(1.0), y, Inches(2.8), Inches(0.35), t['phase'], sz=13, bold=True)
            self.tb(s, Inches(4.0), y, Inches(5.5), Inches(0.35), t['content'], sz=12, color=self.GRAY)
            self.tb(s, Inches(10), y, Inches(2), Inches(0.35), t['duration'], sz=13, align=PP_ALIGN.RIGHT)

        ts = self.a.get('timeline_summary', '')
        if ts:
            ty = y0 + Inches(0.45) + Inches(len(self.a.get('timeline', [])) * 0.42) + Inches(0.2)
            self.tb(s, Inches(1.0), ty, Inches(10), Inches(0.3), ts, sz=12, color=self.LGRAY)

    def build_terms(self):
        s = self.new_slide()
        # Payment
        self.tb(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5), '付款方式', sz=22, bold=True)
        self.line(s, Inches(0.8), Inches(0.95), Inches(11.5))
        self.bullets(s, Inches(1.0), Inches(1.2), Inches(10), Inches(1.5),
                     self.a.get('payment_terms', []), sz=15, sp=Pt(10))

        note = self.a.get('payment_note', '')
        if note:
            self.tb(s, Inches(1.0), Inches(2.5), Inches(10), Inches(0.3), note, sz=11, color=self.LGRAY)

        # Service terms
        self.tb(s, Inches(0.8), Inches(3.0), Inches(5), Inches(0.5), '服務條款', sz=22, bold=True)
        self.line(s, Inches(0.8), Inches(3.55), Inches(11.5))
        self.bullets(s, Inches(1.0), Inches(3.8), Inches(11), Inches(3.5),
                     self.a.get('service_terms', []), sz=13, sp=Pt(6))

    def build_why_us(self):
        s = self.new_slide()
        self.slide_title(s, '為什麼選擇我們', 8)

        why_items = []
        for w in self.a.get('why_us', []):
            why_items.append(f"{w['title']}: {w['description']}")
        if why_items:
            self.bullets(s, Inches(1.0), Inches(1.3), Inches(11), Inches(2.5),
                         why_items, sz=15, sp=Pt(14))

        # Contact
        company = self.c.get('company', {})
        if company.get('name') or company.get('website'):
            self.tb(s, Inches(0.8), Inches(3.8), Inches(8), Inches(0.5), '聯絡資訊', sz=22, bold=True)
            self.line(s, Inches(0.8), Inches(4.35), Inches(11.5))
            self.box(s, Inches(0.8), Inches(4.6), Inches(11.5), Inches(1.5))

            y = Inches(4.8)
            if company.get('name'):
                self.tb(s, Inches(1.2), y, Inches(4), Inches(0.4), company['name'], sz=18, bold=True)
                y += Inches(0.5)
            if company.get('title'):
                self.tb(s, Inches(1.2), y, Inches(4), Inches(0.35), company['title'], sz=14, color=self.GRAY)
                y += Inches(0.4)
            if company.get('website'):
                self.tb(s, Inches(1.2), y, Inches(4), Inches(0.35),
                        f"Website: {company['website']}", sz=14, color=self.GRAY)

    def generate(self):
        self.build_cover()
        self.build_background()
        self.build_architecture()
        self.build_features()

        plans = self.a.get('pricing', {}).get('plans', [])
        for i, plan in enumerate(plans):
            self.build_pricing(plan, slide_num=4 + i, is_recommended=plan.get('recommended', False))

        self.build_pricing_comparison()
        self.build_maintenance_timeline()
        self.build_gantt_timeline()
        self.build_terms()
        self.build_why_us()

        self.prs.save(self.output_path)


def main():
    data = json.loads(sys.stdin.read())
    gen = ProposalGenerator(data['analysis'], data['config'], data['outputPath'])
    gen.generate()


if __name__ == '__main__':
    main()
